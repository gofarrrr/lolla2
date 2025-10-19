#!/usr/bin/env python3
"""
METIS Professional Deliverable Generator
Creates McKinsey-grade client deliverables in multiple formats with brutalist styling
"""

import logging
from datetime import datetime
from typing import Dict, List, Optional, Any
from uuid import UUID
from pathlib import Path
import json
from dataclasses import dataclass, asdict

try:
    # PDF generation
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        Table,
        TableStyle,
        PageBreak,
    )
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch, mm
    from reportlab.pdfgen import canvas
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY

    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print(
        "Warning: PDF generation not available. Install reportlab: pip install reportlab"
    )

try:
    # PowerPoint generation
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print(
        "Warning: PowerPoint generation not available. Install python-pptx: pip install python-pptx"
    )

try:
    # Word generation
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print(
        "Warning: Word generation not available. Install python-docx: pip install python-docx"
    )


# Brutalist Design Constants
BRUTALIST_COLORS = {
    "black": (0, 0, 0),
    "white": (255, 255, 255),
    "neon": (0, 255, 65),
    "gray_light": (245, 245, 245),
    "gray_medium": (204, 204, 204),
}

BRUTALIST_FONTS = {
    "title": ("IBM-Plex-Mono", 28),
    "header": ("IBM-Plex-Mono", 20),
    "subheader": ("IBM-Plex-Mono", 16),
    "body": ("IBM-Plex-Mono", 12),
    "small": ("IBM-Plex-Mono", 10),
}


@dataclass
class DeliverableContent:
    """Structured content for deliverable generation"""

    engagement_id: UUID
    client_name: str
    title: str
    subtitle: str
    executive_summary: Dict[str, Any]
    key_findings: List[Dict[str, Any]]
    recommendations: List[Dict[str, Any]]
    implementation_roadmap: List[Dict[str, Any]]
    appendix: List[Dict[str, Any]]
    metadata: Dict[str, Any]
    generated_at: datetime = datetime.utcnow()


@dataclass
class ExportConfiguration:
    """Configuration for deliverable export"""

    format: str = "pdf"  # pdf, pptx, docx, json
    template: str = "mckinsey"  # mckinsey, bcg, bain, brutalist
    include_appendix: bool = True
    include_watermark: bool = True
    page_numbers: bool = True
    table_of_contents: bool = True
    brutalist_styling: bool = True


class BrutalistDeliverableGenerator:
    """
    Generates professional consulting deliverables with brutalist design
    Supports PDF, PowerPoint, Word, and JSON formats
    """

    def __init__(self, output_dir: Optional[Path] = None):
        self.output_dir = output_dir or Path("./exports")
        self.output_dir.mkdir(exist_ok=True)
        self.logger = logging.getLogger(__name__)

        # Templates directory
        self.templates_dir = Path(__file__).parent / "templates"
        self.templates_dir.mkdir(exist_ok=True)

    async def generate_deliverable(
        self, content: DeliverableContent, config: ExportConfiguration
    ) -> Dict[str, Any]:
        """Generate deliverable in specified format"""

        start_time = datetime.utcnow()

        try:
            if config.format == "pdf":
                file_path = await self._generate_pdf(content, config)
            elif config.format == "pptx":
                file_path = await self._generate_powerpoint(content, config)
            elif config.format == "docx":
                file_path = await self._generate_word(content, config)
            elif config.format == "json":
                file_path = await self._generate_json(content, config)
            else:
                raise ValueError(f"Unsupported format: {config.format}")

            generation_time = (datetime.utcnow() - start_time).total_seconds()

            return {
                "status": "success",
                "file_path": str(file_path),
                "format": config.format,
                "template": config.template,
                "generation_time": generation_time,
                "file_size": file_path.stat().st_size if file_path.exists() else 0,
                "pages": self._estimate_page_count(content, config.format),
                "generated_at": datetime.utcnow().isoformat(),
            }

        except Exception as e:
            self.logger.error(f"Deliverable generation failed: {str(e)}")
            return {
                "status": "error",
                "error": str(e),
                "generation_time": (datetime.utcnow() - start_time).total_seconds(),
            }

    async def _generate_pdf(
        self, content: DeliverableContent, config: ExportConfiguration
    ) -> Path:
        """Generate PDF deliverable with brutalist styling"""

        if not PDF_AVAILABLE:
            raise RuntimeError("PDF generation not available. Install reportlab.")

        filename = f"metis_deliverable_{content.engagement_id}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pdf"
        file_path = self.output_dir / filename

        # Create PDF document
        doc = SimpleDocTemplate(
            str(file_path),
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72,
        )

        # Custom brutalist styles
        styles = self._create_brutalist_pdf_styles()
        story = []

        # Title page
        story.extend(self._create_pdf_title_page(content, styles))
        story.append(PageBreak())

        # Table of contents
        if config.table_of_contents:
            story.extend(self._create_pdf_table_of_contents(content, styles))
            story.append(PageBreak())

        # Executive summary
        story.extend(self._create_pdf_executive_summary(content, styles))
        story.append(PageBreak())

        # Key findings
        story.extend(self._create_pdf_key_findings(content, styles))
        story.append(PageBreak())

        # Recommendations
        story.extend(self._create_pdf_recommendations(content, styles))
        story.append(PageBreak())

        # Implementation roadmap
        story.extend(self._create_pdf_implementation(content, styles))

        # Appendix
        if config.include_appendix and content.appendix:
            story.append(PageBreak())
            story.extend(self._create_pdf_appendix(content, styles))

        # Build PDF
        doc.build(story)

        self.logger.info(f"PDF deliverable generated: {file_path}")
        return file_path

    async def _generate_powerpoint(
        self, content: DeliverableContent, config: ExportConfiguration
    ) -> Path:
        """Generate PowerPoint deliverable with brutalist styling"""

        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "PowerPoint generation not available. Install python-pptx."
            )

        filename = f"metis_presentation_{content.engagement_id}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pptx"
        file_path = self.output_dir / filename

        # Create presentation
        prs = Presentation()

        # Apply brutalist theme
        self._apply_brutalist_pptx_theme(prs)

        # Title slide
        self._create_pptx_title_slide(prs, content)

        # Agenda slide
        self._create_pptx_agenda_slide(prs, content)

        # Executive summary slide
        self._create_pptx_executive_summary_slide(prs, content)

        # Key findings slides
        for finding in content.key_findings:
            self._create_pptx_finding_slide(prs, finding)

        # Recommendations slides
        for i, recommendation in enumerate(content.recommendations):
            self._create_pptx_recommendation_slide(prs, recommendation, i + 1)

        # Implementation roadmap slide
        self._create_pptx_roadmap_slide(prs, content)

        # Appendix slides
        if config.include_appendix:
            for appendix_item in content.appendix:
                self._create_pptx_appendix_slide(prs, appendix_item)

        # Save presentation
        prs.save(str(file_path))

        self.logger.info(f"PowerPoint deliverable generated: {file_path}")
        return file_path

    async def _generate_word(
        self, content: DeliverableContent, config: ExportConfiguration
    ) -> Path:
        """Generate Word deliverable with brutalist styling"""

        if not DOCX_AVAILABLE:
            raise RuntimeError("Word generation not available. Install python-docx.")

        filename = f"metis_report_{content.engagement_id}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.docx"
        file_path = self.output_dir / filename

        # Create document
        doc = Document()

        # Apply brutalist styling
        self._apply_brutalist_docx_styling(doc)

        # Title page
        self._create_docx_title_page(doc, content)

        # Table of contents placeholder
        if config.table_of_contents:
            doc.add_page_break()
            self._create_docx_toc(doc)

        # Executive summary
        doc.add_page_break()
        self._create_docx_executive_summary(doc, content)

        # Key findings
        doc.add_page_break()
        self._create_docx_key_findings(doc, content)

        # Recommendations
        doc.add_page_break()
        self._create_docx_recommendations(doc, content)

        # Implementation roadmap
        doc.add_page_break()
        self._create_docx_implementation(doc, content)

        # Appendix
        if config.include_appendix and content.appendix:
            doc.add_page_break()
            self._create_docx_appendix(doc, content)

        # Save document
        doc.save(str(file_path))

        self.logger.info(f"Word deliverable generated: {file_path}")
        return file_path

    async def _generate_json(
        self, content: DeliverableContent, config: ExportConfiguration
    ) -> Path:
        """Generate JSON deliverable for API consumption"""

        filename = f"metis_data_{content.engagement_id}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.json"
        file_path = self.output_dir / filename

        # Create structured JSON output
        deliverable_data = {
            "metis_deliverable": {
                "version": "7.0.0",
                "generated_at": datetime.utcnow().isoformat(),
                "engagement": {
                    "id": str(content.engagement_id),
                    "client": content.client_name,
                    "title": content.title,
                    "subtitle": content.subtitle,
                },
                "content": {
                    "executive_summary": content.executive_summary,
                    "key_findings": content.key_findings,
                    "recommendations": content.recommendations,
                    "implementation_roadmap": content.implementation_roadmap,
                    "appendix": content.appendix if config.include_appendix else [],
                },
                "metadata": content.metadata,
                "export_config": asdict(config),
            }
        }

        # Write to file
        with open(file_path, "w", encoding="utf-8") as f:
            json.dump(deliverable_data, f, indent=2, ensure_ascii=False, default=str)

        self.logger.info(f"JSON deliverable generated: {file_path}")
        return file_path

    def _create_brutalist_pdf_styles(self):
        """Create brutalist PDF styles"""
        styles = getSampleStyleSheet()

        # Title style
        styles.add(
            ParagraphStyle(
                name="BrutalistTitle",
                parent=styles["Title"],
                fontName="Helvetica-Bold",
                fontSize=36,
                textColor=colors.black,
                alignment=TA_LEFT,
                spaceAfter=30,
                spaceBefore=0,
            )
        )

        # Header style
        styles.add(
            ParagraphStyle(
                name="BrutalistHeader",
                parent=styles["Heading1"],
                fontName="Helvetica-Bold",
                fontSize=20,
                textColor=colors.black,
                alignment=TA_LEFT,
                spaceAfter=20,
                spaceBefore=20,
                borderWidth=4,
                borderColor=colors.black,
                borderPadding=10,
                backColor=colors.Color(0, 1, 0.254),  # Neon green
            )
        )

        # Subheader style
        styles.add(
            ParagraphStyle(
                name="BrutalistSubheader",
                parent=styles["Heading2"],
                fontName="Helvetica-Bold",
                fontSize=16,
                textColor=colors.black,
                alignment=TA_LEFT,
                spaceAfter=12,
                spaceBefore=12,
            )
        )

        # Body style
        styles.add(
            ParagraphStyle(
                name="BrutalistBody",
                parent=styles["Normal"],
                fontName="Helvetica",
                fontSize=12,
                textColor=colors.black,
                alignment=TA_JUSTIFY,
                spaceAfter=12,
                leading=16,
            )
        )

        # Insight box style
        styles.add(
            ParagraphStyle(
                name="BrutalistInsight",
                parent=styles["Normal"],
                fontName="Helvetica-Bold",
                fontSize=12,
                textColor=colors.black,
                alignment=TA_LEFT,
                borderWidth=4,
                borderColor=colors.black,
                borderPadding=15,
                backColor=colors.Color(0, 1, 0.254),  # Neon green
                spaceAfter=15,
                spaceBefore=15,
            )
        )

        return styles

    def _create_pdf_title_page(self, content: DeliverableContent, styles):
        """Create brutalist PDF title page"""
        story = []

        # Main title
        title = f"<b>{content.title.upper()}</b>"
        story.append(Paragraph(title, styles["BrutalistTitle"]))
        story.append(Spacer(1, 20))

        # Subtitle
        subtitle = f"<b>{content.subtitle.upper()}</b>"
        story.append(Paragraph(subtitle, styles["BrutalistHeader"]))
        story.append(Spacer(1, 40))

        # Client info
        client_info = f"<b>CLIENT: {content.client_name.upper()}</b>"
        story.append(Paragraph(client_info, styles["BrutalistSubheader"]))
        story.append(Spacer(1, 20))

        # Metis branding
        branding = "<b>METIS COGNITIVE INTELLIGENCE PLATFORM</b>"
        story.append(Paragraph(branding, styles["BrutalistBody"]))
        story.append(Spacer(1, 10))

        # Date
        date = f"<b>GENERATED: {content.generated_at.strftime('%d %B %Y').upper()}</b>"
        story.append(Paragraph(date, styles["BrutalistBody"]))

        return story

    def _create_pdf_executive_summary(self, content: DeliverableContent, styles):
        """Create executive summary section"""
        story = []

        story.append(Paragraph("<b>EXECUTIVE SUMMARY</b>", styles["BrutalistHeader"]))

        summary = content.executive_summary

        # Governing thought
        if "governing_thought" in summary:
            story.append(
                Paragraph(
                    f"<b>KEY INSIGHT:</b> {summary['governing_thought']}",
                    styles["BrutalistInsight"],
                )
            )

        # Key recommendations
        if "key_recommendations" in summary:
            story.append(
                Paragraph(
                    "<b>STRATEGIC RECOMMENDATIONS</b>", styles["BrutalistSubheader"]
                )
            )
            for i, rec in enumerate(summary["key_recommendations"], 1):
                story.append(Paragraph(f"<b>{i}.</b> {rec}", styles["BrutalistBody"]))
                story.append(Spacer(1, 8))

        # Financial impact
        if "investment_required" in summary and "implementation_timeline" in summary:
            story.append(Spacer(1, 20))
            financial_info = f"""
            <b>INVESTMENT:</b> {summary.get('investment_required', 'TBD')}<br/>
            <b>TIMELINE:</b> {summary.get('implementation_timeline', 'TBD')}<br/>
            <b>EXPECTED ROI:</b> {summary.get('expected_roi', 'High')}
            """
            story.append(Paragraph(financial_info, styles["BrutalistInsight"]))

        return story

    def _create_pdf_key_findings(self, content: DeliverableContent, styles):
        """Create key findings section"""
        story = []

        story.append(Paragraph("<b>KEY FINDINGS</b>", styles["BrutalistHeader"]))

        for i, finding in enumerate(content.key_findings, 1):
            # Finding title
            title = finding.get("title", f"Finding {i}")
            story.append(
                Paragraph(f"<b>{i}. {title.upper()}</b>", styles["BrutalistSubheader"])
            )

            # Finding description
            description = finding.get("description", "")
            story.append(Paragraph(description, styles["BrutalistBody"]))

            # Evidence/Impact
            if "evidence" in finding:
                evidence = f"<b>EVIDENCE:</b> {finding['evidence']}"
                story.append(Paragraph(evidence, styles["BrutalistInsight"]))

            story.append(Spacer(1, 15))

        return story

    def _create_pdf_recommendations(self, content: DeliverableContent, styles):
        """Create recommendations section"""
        story = []

        story.append(
            Paragraph("<b>STRATEGIC RECOMMENDATIONS</b>", styles["BrutalistHeader"])
        )

        for i, rec in enumerate(content.recommendations, 1):
            # Recommendation title
            title = rec.get("title", f"Recommendation {i}")
            story.append(
                Paragraph(f"<b>{i}. {title.upper()}</b>", styles["BrutalistSubheader"])
            )

            # Description
            description = rec.get("description", "")
            story.append(Paragraph(description, styles["BrutalistBody"]))

            # Impact and timeline
            impact_info = []
            if "expected_impact" in rec:
                impact_info.append(f"<b>IMPACT:</b> {rec['expected_impact']}")
            if "timeline" in rec:
                impact_info.append(f"<b>TIMELINE:</b> {rec['timeline']}")
            if "priority" in rec:
                impact_info.append(f"<b>PRIORITY:</b> {rec['priority'].upper()}")

            if impact_info:
                story.append(
                    Paragraph("<br/>".join(impact_info), styles["BrutalistInsight"])
                )

            story.append(Spacer(1, 15))

        return story

    def _create_pdf_implementation(self, content: DeliverableContent, styles):
        """Create implementation roadmap section"""
        story = []

        story.append(
            Paragraph("<b>IMPLEMENTATION ROADMAP</b>", styles["BrutalistHeader"])
        )

        # Create roadmap table
        if content.implementation_roadmap:
            table_data = [["PHASE", "TIMELINE", "KEY ACTIVITIES", "DELIVERABLES"]]

            for phase in content.implementation_roadmap:
                table_data.append(
                    [
                        phase.get("phase", ""),
                        phase.get("timeline", ""),
                        phase.get("activities", ""),
                        phase.get("deliverables", ""),
                    ]
                )

            table = Table(
                table_data, colWidths=[1.5 * inch, 1 * inch, 2.5 * inch, 2 * inch]
            )
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.black),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("FONTSIZE", (0, 0), (-1, 0), 12),
                        ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                        (
                            "BACKGROUND",
                            (0, 1),
                            (-1, -1),
                            colors.Color(245 / 255, 245 / 255, 245 / 255),
                        ),
                        ("GRID", (0, 0), (-1, -1), 2, colors.black),
                    ]
                )
            )

            story.append(table)

        return story

    def _create_pdf_appendix(self, content: DeliverableContent, styles):
        """Create appendix section"""
        story = []

        story.append(Paragraph("<b>APPENDIX</b>", styles["BrutalistHeader"]))

        for item in content.appendix:
            title = item.get("title", "Appendix Item")
            story.append(
                Paragraph(f"<b>{title.upper()}</b>", styles["BrutalistSubheader"])
            )

            content_text = item.get("content", "")
            story.append(Paragraph(content_text, styles["BrutalistBody"]))

            story.append(Spacer(1, 20))

        return story

    def _create_pdf_table_of_contents(self, content: DeliverableContent, styles):
        """Create table of contents"""
        story = []

        story.append(Paragraph("<b>TABLE OF CONTENTS</b>", styles["BrutalistHeader"]))

        toc_items = [
            "1. EXECUTIVE SUMMARY",
            "2. KEY FINDINGS",
            "3. STRATEGIC RECOMMENDATIONS",
            "4. IMPLEMENTATION ROADMAP",
        ]

        if content.appendix:
            toc_items.append("5. APPENDIX")

        for item in toc_items:
            story.append(Paragraph(item, styles["BrutalistBody"]))
            story.append(Spacer(1, 8))

        return story

    # PowerPoint helper methods (simplified for brevity)
    def _apply_brutalist_pptx_theme(self, prs):
        """Apply brutalist theme to PowerPoint"""
        # This would customize the master slides with brutalist styling
        pass

    def _create_pptx_title_slide(self, prs, content):
        """Create PowerPoint title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = content.title.upper()
        subtitle.text = f"{content.subtitle.upper()}\n{content.client_name.upper()}"

        # Apply brutalist styling
        title.text_frame.paragraphs[0].font.name = "IBM Plex Mono"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True

    def _create_pptx_agenda_slide(self, prs, content):
        """Create agenda slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = "AGENDA"

        agenda_items = [
            "EXECUTIVE SUMMARY",
            "KEY FINDINGS",
            "STRATEGIC RECOMMENDATIONS",
            "IMPLEMENTATION ROADMAP",
        ]

        tf = content_placeholder.text_frame
        for item in agenda_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.name = "IBM Plex Mono"
            p.font.size = Pt(20)
            p.font.bold = True

    def _create_pptx_executive_summary_slide(self, prs, content):
        """Create executive summary slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = "EXECUTIVE SUMMARY"

        summary = content.executive_summary
        tf = content_placeholder.text_frame

        if "governing_thought" in summary:
            p = tf.paragraphs[0]
            p.text = f"KEY INSIGHT: {summary['governing_thought']}"
            p.font.name = "IBM Plex Mono"
            p.font.size = Pt(18)
            p.font.bold = True

    def _create_pptx_finding_slide(self, prs, finding):
        """Create finding slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = finding.get("title", "KEY FINDING").upper()

        tf = content_placeholder.text_frame
        tf.text = finding.get("description", "")

    def _create_pptx_recommendation_slide(self, prs, recommendation, number):
        """Create recommendation slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = (
            f"RECOMMENDATION {number}: {recommendation.get('title', '').upper()}"
        )

        tf = content_placeholder.text_frame
        tf.text = recommendation.get("description", "")

    def _create_pptx_roadmap_slide(self, prs, content):
        """Create roadmap slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "IMPLEMENTATION ROADMAP"

    def _create_pptx_appendix_slide(self, prs, appendix_item):
        """Create appendix slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = appendix_item.get("title", "APPENDIX").upper()

        tf = content_placeholder.text_frame
        tf.text = appendix_item.get("content", "")

    # Word document helper methods (simplified)
    def _apply_brutalist_docx_styling(self, doc):
        """Apply brutalist styling to Word document"""
        # This would customize document styles
        pass

    def _create_docx_title_page(self, doc, content):
        """Create Word title page"""
        title = doc.add_heading(content.title.upper(), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.LEFT

        subtitle = doc.add_heading(content.subtitle.upper(), 1)
        subtitle.alignment = WD_ALIGN_PARAGRAPH.LEFT

        doc.add_paragraph(f"CLIENT: {content.client_name.upper()}")
        doc.add_paragraph(
            f"GENERATED: {content.generated_at.strftime('%d %B %Y').upper()}"
        )

    def _create_docx_toc(self, doc):
        """Create table of contents placeholder"""
        doc.add_heading("TABLE OF CONTENTS", 1)
        doc.add_paragraph("1. EXECUTIVE SUMMARY")
        doc.add_paragraph("2. KEY FINDINGS")
        doc.add_paragraph("3. STRATEGIC RECOMMENDATIONS")
        doc.add_paragraph("4. IMPLEMENTATION ROADMAP")

    def _create_docx_executive_summary(self, doc, content):
        """Create executive summary"""
        doc.add_heading("EXECUTIVE SUMMARY", 1)

        summary = content.executive_summary
        if "governing_thought" in summary:
            p = doc.add_paragraph()
            p.add_run("KEY INSIGHT: ").bold = True
            p.add_run(summary["governing_thought"])

    def _create_docx_key_findings(self, doc, content):
        """Create key findings section"""
        doc.add_heading("KEY FINDINGS", 1)

        for i, finding in enumerate(content.key_findings, 1):
            doc.add_heading(f"{i}. {finding.get('title', f'Finding {i}').upper()}", 2)
            doc.add_paragraph(finding.get("description", ""))

    def _create_docx_recommendations(self, doc, content):
        """Create recommendations section"""
        doc.add_heading("STRATEGIC RECOMMENDATIONS", 1)

        for i, rec in enumerate(content.recommendations, 1):
            doc.add_heading(
                f"{i}. {rec.get('title', f'Recommendation {i}').upper()}", 2
            )
            doc.add_paragraph(rec.get("description", ""))

    def _create_docx_implementation(self, doc, content):
        """Create implementation section"""
        doc.add_heading("IMPLEMENTATION ROADMAP", 1)

        for phase in content.implementation_roadmap:
            doc.add_heading(phase.get("phase", "").upper(), 2)
            doc.add_paragraph(f"Timeline: {phase.get('timeline', '')}")
            doc.add_paragraph(f"Activities: {phase.get('activities', '')}")

    def _create_docx_appendix(self, doc, content):
        """Create appendix section"""
        doc.add_heading("APPENDIX", 1)

        for item in content.appendix:
            doc.add_heading(item.get("title", "Appendix Item").upper(), 2)
            doc.add_paragraph(item.get("content", ""))

    def _estimate_page_count(self, content: DeliverableContent, format: str) -> int:
        """Estimate page count based on content"""
        base_pages = 4  # Title, TOC, Executive Summary, Implementation
        base_pages += len(content.key_findings)
        base_pages += len(content.recommendations)
        base_pages += len(content.appendix) if content.appendix else 0

        if format == "pptx":
            return base_pages + 2  # Agenda and closing slide

        return max(base_pages, 8)  # Minimum 8 pages for professional look


# Factory function for easy integration
def create_deliverable_from_engagement(
    engagement_data: Dict[str, Any], synthesis_result: Dict[str, Any]
) -> DeliverableContent:
    """Create DeliverableContent from engagement data"""

    # Extract executive summary
    executive_summary = synthesis_result.get("executive_summary", {})

    # Extract key findings
    key_findings = []
    if "key_findings" in synthesis_result:
        for finding in synthesis_result["key_findings"]:
            key_findings.append(
                {
                    "title": finding.get("title", "Key Finding"),
                    "description": finding.get("description", ""),
                    "evidence": finding.get("evidence", ""),
                    "confidence": finding.get("confidence", 0.8),
                }
            )

    # Extract recommendations
    recommendations = []
    if "recommendations" in executive_summary:
        for i, rec in enumerate(executive_summary["recommendations"], 1):
            recommendations.append(
                {
                    "title": f"Strategic Recommendation {i}",
                    "description": rec,
                    "priority": "High",
                    "timeline": f"Phase {i}",
                    "expected_impact": "Significant",
                }
            )

    # Create implementation roadmap
    implementation_roadmap = [
        {
            "phase": "Phase 1: Foundation",
            "timeline": "0-3 months",
            "activities": "Infrastructure setup, team formation",
            "deliverables": "Project charter, team structure",
        },
        {
            "phase": "Phase 2: Implementation",
            "timeline": "3-12 months",
            "activities": "Execute key recommendations",
            "deliverables": "Progress reports, milestone achievements",
        },
        {
            "phase": "Phase 3: Optimization",
            "timeline": "12-18 months",
            "activities": "Monitor results, continuous improvement",
            "deliverables": "Final assessment, lessons learned",
        },
    ]

    return DeliverableContent(
        engagement_id=UUID(engagement_data["engagement_id"]),
        client_name=engagement_data["client_name"],
        title=f"{engagement_data['client_name']} Strategic Analysis",
        subtitle="Digital Transformation Strategy",
        executive_summary=executive_summary,
        key_findings=key_findings,
        recommendations=recommendations,
        implementation_roadmap=implementation_roadmap,
        appendix=[],
        metadata={
            "confidence_score": synthesis_result.get("overall_confidence", 0.85),
            "pyramid_compliance": synthesis_result.get("pyramid_compliance", 0.90),
            "partner_ready_score": synthesis_result.get("partner_ready_score", 0.75),
        },
    )
